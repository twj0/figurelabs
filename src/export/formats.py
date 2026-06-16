"""Format conversion and download — all supported output formats in one place.

HAR-confirmed facts (2026-06-15):
  - /api/file-proxy returns Content-Type: image/jpeg regardless of the S3 path suffix.
    PNG and JPG are the same bytes; the extension is chosen client-side.
  - /app-api/plot/image/svg converts the image server-side and returns an S3 SVG URL.
  - PPTX has no server-side endpoint. The web UI generates it client-side from the SVG
    (browser JS, likely pptxgenjs). We replicate this locally with python-pptx.
"""

import os
from typing import Optional

import requests

BASE_URL = "https://chat.figurelabs.ai"
_PROXY       = f"{BASE_URL}/api/file-proxy"
_SVG_CONVERT = f"{BASE_URL}/app-api/plot/image/svg"


# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------

def _stem_from_url(url: str) -> str:
    part = url.split("/")[-1].split("?")[0]
    return part.rsplit(".", 1)[0] or "figure"


def _write_stream(resp: requests.Response, path: str) -> str:
    with open(path, "wb") as f:
        for chunk in resp.iter_content(8192):
            f.write(chunk)
    return path


def _fetch_via_proxy(
    session: requests.Session,
    s3_url: str,
    ext: str,
    output_dir: str,
    filename: Optional[str],
) -> Optional[str]:
    """Download any S3 resource through the file-proxy and save with the given extension."""
    resp = session.get(_PROXY, params={"url": s3_url}, stream=True)
    if resp.status_code != 200:
        print(f"[{ext}] proxy HTTP {resp.status_code}")
        return None
    os.makedirs(output_dir, exist_ok=True)
    path = os.path.join(output_dir, f"{filename or _stem_from_url(s3_url)}.{ext}")
    _write_stream(resp, path)
    print(f"[{ext}] {path} ({os.path.getsize(path) // 1024} KB)")
    return path


def _request_svg_url(session: requests.Session, png_s3_url: str) -> Optional[str]:
    """POST to the SVG conversion endpoint and return the resulting S3 URL."""
    resp = session.post(_SVG_CONVERT, json={"imageUrl": [png_s3_url]})
    if resp.status_code != 200:
        print(f"[svg] conversion HTTP {resp.status_code}")
        return None
    result = resp.json()
    if result.get("code") != 0:
        print(f"[svg] conversion failed: {result.get('msg', result)}")
        return None
    urls = result["data"].get("fileUrls", [])
    if not urls:
        print("[svg] no fileUrls in response")
        return None
    return urls[0]


# ---------------------------------------------------------------------------
# Public download functions
# ---------------------------------------------------------------------------

def download_png(
    session: requests.Session,
    png_s3_url: str,
    output_dir: str = ".",
    filename: Optional[str] = None,
) -> Optional[str]:
    # Server returns image/jpeg regardless of S3 path suffix — saved as .png for compatibility.
    return _fetch_via_proxy(session, png_s3_url, "png", output_dir, filename)


def download_jpg(
    session: requests.Session,
    png_s3_url: str,
    output_dir: str = ".",
    filename: Optional[str] = None,
) -> Optional[str]:
    # Identical bytes to PNG download; only the saved extension differs.
    return _fetch_via_proxy(session, png_s3_url, "jpg", output_dir, filename)


def download_svg(
    session: requests.Session,
    png_s3_url: str,
    output_dir: str = ".",
    filename: Optional[str] = None,
) -> Optional[str]:
    svg_url = _request_svg_url(session, png_s3_url)
    if not svg_url:
        return None
    return _fetch_via_proxy(session, svg_url, "svg", output_dir, filename)


def download_pptx(
    session: requests.Session,
    png_s3_url: str,
    output_dir: str = ".",
    filename: Optional[str] = None,
) -> Optional[str]:
    """Build a PPTX with a proper OOXML SVG embed (PowerPoint 2016+ vector-editable).

    OOXML SVG extension (ECMA-376 §14.2.1 + MS-PPTX §2.3.2):
      <a:blip r:embed="rId1">          ← PNG fallback (required for older readers)
        <a:extLst>
          <a:ext uri="{96DAC541-7B7A-43D3-8B79-37D633B846F1}">
            <asvg:svgBlip r:embed="rId2"/>  ← actual SVG (vector-quality in PPT 2016+)
          </a:ext>
        </a:extLst>
      </a:blip>

    Without this structure PowerPoint treats the embedded image as a raster blip
    and rasterises it at screen resolution — the bug in the previous implementation.
    """
    try:
        from pptx import Presentation
        from pptx.util import Emu
        from pptx.oxml.ns import qn
        from lxml import etree
        import io, zipfile
    except ImportError:
        print("[pptx] python-pptx or lxml not installed — run: uv pip install python-pptx lxml")
        return None

    # 1. Fetch PNG bytes (fallback blip for older PowerPoint readers)
    png_resp = session.get(_PROXY, params={"url": png_s3_url}, stream=True)
    if png_resp.status_code != 200:
        print(f"[pptx] PNG fetch HTTP {png_resp.status_code}")
        return None
    png_bytes = png_resp.content

    # 2. Fetch SVG bytes
    svg_url = _request_svg_url(session, png_s3_url)
    if not svg_url:
        return None
    svg_resp = session.get(_PROXY, params={"url": svg_url}, stream=True)
    if svg_resp.status_code != 200:
        print(f"[pptx] SVG fetch HTTP {svg_resp.status_code}")
        return None
    svg_bytes = svg_resp.content

    # 3. Parse SVG viewBox for slide aspect ratio
    try:
        root = etree.fromstring(svg_bytes)
        vb = root.get("viewBox", "")
        parts = vb.split()
        if len(parts) == 4:
            svg_w, svg_h = float(parts[2]), float(parts[3])
        else:
            w_attr = root.get("width", "1376")
            h_attr = root.get("height", "768")
            svg_w = float("".join(c for c in w_attr if c.isdigit() or c == ".") or 1376)
            svg_h = float("".join(c for c in h_attr if c.isdigit() or c == ".") or 768)
    except Exception:
        svg_w, svg_h = 1376.0, 768.0

    # 4. Build blank PPTX in memory
    prs = Presentation()
    slide_w = Emu(9144000)  # 10 inches
    slide_h = Emu(int(slide_w * svg_h / svg_w))
    prs.slide_width = slide_w
    prs.slide_height = slide_h
    prs.slides.add_slide(prs.slide_layouts[6])  # blank

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    # 5. Patch PPTX ZIP: inject PNG + SVG, fix slide XML / rels / content-types
    out_buf = io.BytesIO()
    slide_xml_name = "ppt/slides/slide1.xml"
    rels_name = "ppt/slides/_rels/slide1.xml.rels"
    ct_name = "[Content_Types].xml"
    rewritten = {slide_xml_name, rels_name, ct_name}

    _REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
    _IMG_TYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
    _A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
    _P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
    _R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    _ASVG_NS = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"
    _SVG_EXT_URI = "{96DAC541-7B7A-43D3-8B79-37D633B846F1}"

    with zipfile.ZipFile(buf, "r") as zin, zipfile.ZipFile(out_buf, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            if item.filename not in rewritten:
                zout.writestr(item, zin.read(item.filename))

        zout.writestr("ppt/media/image1.png", png_bytes)
        zout.writestr("ppt/media/image2.svg", svg_bytes)

        # --- slide1.xml: add <p:pic> with PNG blip + asvg:svgBlip extension ---
        slide_xml = zin.read(slide_xml_name)
        tree = etree.fromstring(slide_xml)
        spTree = tree.find(".//" + qn("p:spTree"))

        pic_xml = (
            f'<p:pic'
            f' xmlns:p="{_P_NS}"'
            f' xmlns:a="{_A_NS}"'
            f' xmlns:r="{_R_NS}"'
            f' xmlns:asvg="{_ASVG_NS}">'
            f'<p:nvPicPr>'
            f'<p:cNvPr id="2" name="Figure"/>'
            f'<p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>'
            f'<p:nvPr/>'
            f'</p:nvPicPr>'
            f'<p:blipFill>'
            f'<a:blip r:embed="rId1">'
            f'<a:extLst>'
            f'<a:ext uri="{_SVG_EXT_URI}">'
            f'<asvg:svgBlip r:embed="rId2"/>'
            f'</a:ext>'
            f'</a:extLst>'
            f'</a:blip>'
            f'<a:stretch><a:fillRect/></a:stretch>'
            f'</p:blipFill>'
            f'<p:spPr>'
            f'<a:xfrm><a:off x="0" y="0"/><a:ext cx="{int(slide_w)}" cy="{int(slide_h)}"/></a:xfrm>'
            f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
            f'</p:spPr>'
            f'</p:pic>'
        )
        spTree.append(etree.fromstring(pic_xml))
        zout.writestr(slide_xml_name, etree.tostring(tree, xml_declaration=True,
                                                      encoding="UTF-8", standalone=True))

        # --- slide1.xml.rels: rId1=PNG, rId2=SVG ---
        try:
            rels_tree = etree.fromstring(zin.read(rels_name))
        except KeyError:
            rels_tree = etree.fromstring(
                b'<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"/>'
            )
        for rid, target in [("rId1", "../media/image1.png"), ("rId2", "../media/image2.svg")]:
            rel = etree.SubElement(rels_tree, f"{{{_REL_NS}}}Relationship")
            rel.set("Id", rid)
            rel.set("Type", _IMG_TYPE)
            rel.set("Target", target)
        zout.writestr(rels_name, etree.tostring(rels_tree, xml_declaration=True,
                                                encoding="UTF-8", standalone=True))

        # --- [Content_Types].xml: register svg and png extensions ---
        ct_xml = zin.read(ct_name)
        ct_tree = etree.fromstring(ct_xml)
        ct_ns = "http://schemas.openxmlformats.org/package/2006/content-types"
        existing = {e.get("Extension") for e in ct_tree.findall(f"{{{ct_ns}}}Default")}
        for ext, ct in [("svg", "image/svg+xml"), ("png", "image/png")]:
            if ext not in existing:
                el = etree.SubElement(ct_tree, f"{{{ct_ns}}}Default")
                el.set("Extension", ext)
                el.set("ContentType", ct)
        zout.writestr(ct_name, etree.tostring(ct_tree, xml_declaration=True,
                                              encoding="UTF-8", standalone=True))

    os.makedirs(output_dir, exist_ok=True)
    stem = filename or _stem_from_url(png_s3_url)
    path = os.path.join(output_dir, f"{stem}.pptx")
    with open(path, "wb") as f:
        f.write(out_buf.getvalue())

    print(f"[pptx] {path} ({os.path.getsize(path) // 1024} KB)")
    return path


# ---------------------------------------------------------------------------
# Dispatch table
# ---------------------------------------------------------------------------

_HANDLERS = {
    "png":  download_png,
    "jpg":  download_jpg,
    "jpeg": download_jpg,
    "svg":  download_svg,
    "pptx": download_pptx,
}


def download(
    session: requests.Session,
    png_s3_url: str,
    fmt: str,
    output_dir: str = ".",
    filename: Optional[str] = None,
) -> Optional[str]:
    """Unified entry point — dispatch to the correct format handler."""
    key = fmt.lower().lstrip(".")
    handler = _HANDLERS.get(key)
    if handler is None:
        print(f"[formats] unknown format '{fmt}' — supported: {', '.join(_HANDLERS)}")
        return None
    return handler(session, png_s3_url, output_dir, filename)
