"""FastAPI application — REST API + static frontend hosting."""

import asyncio
import json
import time
from pathlib import Path
from typing import Any, Optional

from fastapi import FastAPI, Header, HTTPException
from fastapi.responses import Response, StreamingResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel

from ..chat.client import FigureLabsChat
from ..config import OPENAI_API_KEY
from ..core.database import stats_db
from ..db import delete_account, init_db, list_accounts, save_account, update_label
from ..export._session import make_session
from ..export.client import ExportClient
from ..export.formats import _PROXY, _request_svg_url
from ..register.client import FigureLabsRegistration
from ..register.mail_service import DuckMailService, MailTmService

app = FastAPI(title="FigureLabs AI", docs_url=None, redoc_url=None)


@app.on_event("startup")
async def startup():
    await init_db()


# ---------------------------------------------------------------------------
# Pydantic models
# ---------------------------------------------------------------------------

class AccountOut(BaseModel):
    id: int
    user_id: str
    email: str
    access_token: str
    mail_service: str
    created_at: int
    label: Optional[str] = None
    expires_time: Optional[int] = None


class RegisterRequest(BaseModel):
    mail_service: str = "mailtm"   # "mailtm" | "duckmail"


class RegisterResponse(BaseModel):
    # For mailtm: account is complete
    # For duckmail: needs manual code entry — returns pending_id
    done: bool
    account: Optional[AccountOut] = None
    # duckmail flow
    pending_id: Optional[str] = None
    email: Optional[str] = None
    inbox_url: Optional[str] = None


class DuckMailVerifyRequest(BaseModel):
    pending_id: str
    code: str


class LabelUpdate(BaseModel):
    label: str


class SessionCreate(BaseModel):
    access_token: str
    title: str = "New Diagram"
    agent_id: int = 0


class SessionResponse(BaseModel):
    session_id: str


class MessageSend(BaseModel):
    access_token: str
    session_id: str
    text: str
    model_id: int = 7
    ratio: str = "16:9"
    scene: str = "gen-svg"
    style: Optional[str] = None


class MessageResponse(BaseModel):
    message_id: str


class StatusResponse(BaseModel):
    status: int
    image_model: Optional[str] = None
    image_ratio: Optional[str] = None
    file_url: list[str] = []


class ExpandRequest(BaseModel):
    access_token: str
    message: str


class ExpandResponse(BaseModel):
    expanded: str


_OPENAI_MODEL_IDS = {
    "figurelabs-nano-banana": 6,
    "figurelabs-nano-banana-pro": 7,
    "figurelabs-banana-2": 10,
    "figurelabs-gpt-image-2": 12,
}
_DEFAULT_OPENAI_MODEL = "figurelabs-nano-banana-pro"


class OpenAIMessage(BaseModel):
    role: str
    content: str | list[dict[str, Any]]


class OpenAIChatRequest(BaseModel):
    model: str = _DEFAULT_OPENAI_MODEL
    messages: list[OpenAIMessage]
    stream: bool = False
    temperature: Optional[float] = None
    top_p: Optional[float] = None
    ratio: str = "16:9"
    scene: str = "gen-svg"
    style: Optional[str] = None
    wait_timeout: int = 180


# In-memory store for pending DuckMail flows {pending_id -> (email, username, code_id, reg)}
_pending_duck: dict[str, tuple] = {}


# ---------------------------------------------------------------------------
# Account management
# ---------------------------------------------------------------------------

@app.get("/api/accounts", response_model=list[AccountOut])
async def api_list_accounts():
    return await list_accounts()


@app.post("/api/accounts/register", response_model=RegisterResponse)
async def api_register(body: RegisterRequest):
    svc = body.mail_service.lower()

    if svc == "mailtm":
        try:
            reg = FigureLabsRegistration(mail_service=MailTmService())
            data = await asyncio.to_thread(reg.register_auto)
            data["mailService"] = "mailtm"
            saved = await save_account(data)
            return RegisterResponse(done=True, account=AccountOut(**saved))
        except Exception as e:
            raise HTTPException(status_code=500, detail=str(e))

    elif svc == "duckmail":
        import uuid as _uuid
        duck = DuckMailService()
        reg = FigureLabsRegistration(mail_service=duck)
        email, username = await asyncio.to_thread(duck.create_account)
        code_id = await asyncio.to_thread(reg.send_verification_code, email)

        pending_id = _uuid.uuid4().hex
        _pending_duck[pending_id] = (email, username, code_id, reg)

        return RegisterResponse(
            done=False,
            pending_id=pending_id,
            email=email,
            inbox_url=f"https://duckmail.sbs/{username}",
        )
    else:
        raise HTTPException(status_code=400, detail="Unknown mail_service")


@app.post("/api/accounts/verify-duckmail", response_model=RegisterResponse)
async def api_verify_duckmail(body: DuckMailVerifyRequest):
    entry = _pending_duck.pop(body.pending_id, None)
    if not entry:
        raise HTTPException(status_code=404, detail="Pending session not found or expired")

    email, username, code_id, reg = entry
    try:
        data = await asyncio.to_thread(reg.login, email, body.code, code_id)
        data["mailService"] = "duckmail"
        data["email"] = email
        saved = await save_account(data)
        return RegisterResponse(done=True, account=AccountOut(**saved))
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))


@app.delete("/api/accounts/{user_id}")
async def api_delete_account(user_id: str):
    ok = await delete_account(user_id)
    if not ok:
        raise HTTPException(status_code=404, detail="Account not found")
    return {"ok": True}


@app.patch("/api/accounts/{user_id}/label")
async def api_update_label(user_id: str, body: LabelUpdate):
    await update_label(user_id, body.label)
    return {"ok": True}


# ---------------------------------------------------------------------------
# Chat
# ---------------------------------------------------------------------------

@app.post("/api/session", response_model=SessionResponse)
async def api_create_session(body: SessionCreate):
    chat = FigureLabsChat(body.access_token)
    sid = await asyncio.to_thread(chat.create_session, body.title, body.agent_id)
    if not sid:
        raise HTTPException(status_code=400, detail="Failed to create session")
    return SessionResponse(session_id=sid)


@app.post("/api/message", response_model=MessageResponse)
async def api_send_message(body: MessageSend):
    chat = FigureLabsChat(body.access_token)
    mid = await asyncio.to_thread(
        chat.send_message,
        body.session_id, body.text,
        "NORMAL_CHAT", body.model_id, body.ratio, body.style, True, body.scene,
    )
    if not mid:
        raise HTTPException(status_code=400, detail="Failed to send message")
    return MessageResponse(message_id=mid)


@app.get("/api/status/{message_id}", response_model=StatusResponse)
async def api_get_status(message_id: str, token: str):
    chat = FigureLabsChat(token)
    st = await asyncio.to_thread(chat.get_message_status, message_id)
    if st is None:
        raise HTTPException(status_code=404, detail="Message not found")
    return StatusResponse(
        status=st.get("status", 0),
        image_model=st.get("imageModel"),
        image_ratio=st.get("imageRatio"),
        file_url=st.get("fileUrl") or [],
    )


@app.post("/api/expand", response_model=ExpandResponse)
async def api_expand(body: ExpandRequest):
    """Expand / enhance a short prompt using the server-side AI."""
    chat = FigureLabsChat(body.access_token)
    expanded = await asyncio.to_thread(chat.expand_prompt, body.message)
    if expanded is None:
        raise HTTPException(status_code=502, detail="Prompt expansion failed")
    return ExpandResponse(expanded=expanded)


# ---------------------------------------------------------------------------
# Download / export
# ---------------------------------------------------------------------------

@app.get("/api/download/{message_id}")
async def api_download(message_id: str, token: str, fmt: str = "png"):
    exp = ExportClient(token)
    urls, file_type = await asyncio.to_thread(exp.get_file_urls, message_id)
    if not urls:
        raise HTTPException(status_code=404, detail="Image not ready")

    file_url = urls[0]
    session = make_session(token)
    fmt = fmt.lower()

    # If source is SVG and user wants SVG, download directly
    if file_type == ".svg" and fmt == "svg":
        resp = await asyncio.to_thread(
            lambda: session.get(file_url, stream=True)
        )
        if resp.status_code != 200:
            raise HTTPException(status_code=502, detail="SVG fetch failed")
        return Response(
            content=resp.content,
            media_type="image/svg+xml",
            headers={"Content-Disposition": 'attachment; filename="figure.svg"'},
        )

    # For PNG outputs or conversions, use the proxy
    if fmt in ("png", "jpg", "jpeg"):
        resp = await asyncio.to_thread(
            lambda: session.get(_PROXY, params={"url": file_url}, stream=True)
        )
        if resp.status_code != 200:
            raise HTTPException(status_code=502, detail="Upstream error")
        media = "image/png" if fmt == "png" else "image/jpeg"
        return Response(
            content=resp.content,
            media_type=media,
            headers={"Content-Disposition": f'attachment; filename="figure.{fmt}"'},
        )

    elif fmt == "svg":
        # Convert PNG to SVG
        svg_url = await asyncio.to_thread(_request_svg_url, session, file_url)
        if not svg_url:
            raise HTTPException(
                status_code=400,
                detail="SVG conversion failed. This may require additional credits or the image may not be convertible to SVG format."
            )
        resp = await asyncio.to_thread(
            lambda: session.get(_PROXY, params={"url": svg_url}, stream=True)
        )
        if resp.status_code != 200:
            raise HTTPException(status_code=502, detail="SVG fetch failed")
        return Response(
            content=resp.content,
            media_type="image/svg+xml",
            headers={"Content-Disposition": 'attachment; filename="figure.svg"'},
        )

    elif fmt == "pptx":
        data = await asyncio.to_thread(_build_pptx_bytes, session, file_url, file_type)
        if data is None:
            raise HTTPException(status_code=502, detail="PPTX build failed")
        return Response(
            content=data,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": 'attachment; filename="figure.pptx"'},
        )

    raise HTTPException(status_code=400, detail=f"Unknown format: {fmt}")


def _build_pptx_bytes(session, file_url: str, file_type: Optional[str] = None) -> Optional[bytes]:
    from ..export.formats import _PROXY, _request_svg_url
    try:
        import io
        import zipfile

        from lxml import etree
        from pptx import Presentation
        from pptx.oxml.ns import qn
        from pptx.util import Emu
    except ImportError:
        return None

    # If source is already SVG, use it directly
    if file_type == ".svg":
        svg_bytes = session.get(file_url, stream=True).content
        # Try to convert SVG to PNG for preview (optional)
        png_bytes = svg_bytes  # Fallback to SVG if no conversion available
    else:
        # Source is PNG
        png_bytes = session.get(_PROXY, params={"url": file_url}, stream=True).content
        svg_url = _request_svg_url(session, file_url)
        if not svg_url:
            return None
        svg_bytes = session.get(_PROXY, params={"url": svg_url}, stream=True).content

    try:
        root = etree.fromstring(svg_bytes)
        parts = root.get("viewBox", "").split()
        svg_w, svg_h = (float(parts[2]), float(parts[3])) if len(parts) == 4 else (1376.0, 768.0)
    except Exception:
        svg_w, svg_h = 1376.0, 768.0

    prs = Presentation()
    slide_w = Emu(9144000)
    slide_h = Emu(int(slide_w * svg_h / svg_w))
    prs.slide_width, prs.slide_height = slide_w, slide_h
    prs.slides.add_slide(prs.slide_layouts[6])

    buf = io.BytesIO(); prs.save(buf); buf.seek(0)

    _REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
    _IMG_TYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
    _A_NS  = "http://schemas.openxmlformats.org/drawingml/2006/main"
    _P_NS  = "http://schemas.openxmlformats.org/presentationml/2006/main"
    _R_NS  = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    _ASVG  = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"
    _EXT   = "{96DAC541-7B7A-43D3-8B79-37D633B846F1}"

    out = io.BytesIO()
    s_xml, rels, ct = "ppt/slides/slide1.xml", "ppt/slides/_rels/slide1.xml.rels", "[Content_Types].xml"
    rewritten = {s_xml, rels, ct}

    with zipfile.ZipFile(buf, "r") as zin, zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            if item.filename not in rewritten:
                zout.writestr(item, zin.read(item.filename))
        zout.writestr("ppt/media/image1.png", png_bytes)
        zout.writestr("ppt/media/image2.svg", svg_bytes)

        tree = etree.fromstring(zin.read(s_xml))
        spTree = tree.find(".//" + qn("p:spTree"))
        pic = (
            f'<p:pic xmlns:p="{_P_NS}" xmlns:a="{_A_NS}" xmlns:r="{_R_NS}" xmlns:asvg="{_ASVG}">'
            f'<p:nvPicPr><p:cNvPr id="2" name="Figure"/>'
            f'<p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr><p:nvPr/></p:nvPicPr>'
            f'<p:blipFill><a:blip r:embed="rId1"><a:extLst><a:ext uri="{_EXT}">'
            f'<asvg:svgBlip r:embed="rId2"/></a:ext></a:extLst></a:blip>'
            f'<a:stretch><a:fillRect/></a:stretch></p:blipFill>'
            f'<p:spPr><a:xfrm><a:off x="0" y="0"/>'
            f'<a:ext cx="{int(slide_w)}" cy="{int(slide_h)}"/></a:xfrm>'
            f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr></p:pic>'
        )
        spTree.append(etree.fromstring(pic))
        zout.writestr(s_xml, etree.tostring(tree, xml_declaration=True, encoding="UTF-8", standalone=True))

        try:
            rt = etree.fromstring(zin.read(rels))
        except KeyError:
            rt = etree.fromstring(b'<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"/>')
        for rid, tgt in [("rId1", "../media/image1.png"), ("rId2", "../media/image2.svg")]:
            r = etree.SubElement(rt, f"{{{_REL_NS}}}Relationship")
            r.set("Id", rid); r.set("Type", _IMG_TYPE); r.set("Target", tgt)
        zout.writestr(rels, etree.tostring(rt, xml_declaration=True, encoding="UTF-8", standalone=True))

        ct_tree = etree.fromstring(zin.read(ct))
        ct_ns = "http://schemas.openxmlformats.org/package/2006/content-types"
        existing = {e.get("Extension") for e in ct_tree.findall(f"{{{ct_ns}}}Default")}
        for ext, ctype in [("svg", "image/svg+xml"), ("png", "image/png")]:
            if ext not in existing:
                el = etree.SubElement(ct_tree, f"{{{ct_ns}}}Default")
                el.set("Extension", ext); el.set("ContentType", ctype)
        zout.writestr(ct, etree.tostring(ct_tree, xml_declaration=True, encoding="UTF-8", standalone=True))

    return out.getvalue()


# ---------------------------------------------------------------------------
# Statistics API (新增)
# ---------------------------------------------------------------------------

class StatsResponse(BaseModel):
    labels: list[str]
    total_requests: list[int]
    failed_requests: list[int]
    rate_limited_requests: list[int]
    model_requests: dict[str, list[int]]
    model_ttfb_times: dict[str, list[float]]
    model_total_times: dict[str, list[float]]


class TotalCountsResponse(BaseModel):
    success: int
    failed: int


@app.get("/api/stats", response_model=StatsResponse)
async def api_get_stats(time_range: str = "24h"):
    """Get aggregated statistics by time range (24h/7d/30d)."""
    if time_range not in ("24h", "7d", "30d"):
        raise HTTPException(status_code=400, detail="Invalid time_range. Use 24h, 7d, or 30d")

    stats = await stats_db.get_stats_by_time_range(time_range)
    return StatsResponse(**stats)


@app.get("/api/stats/totals", response_model=TotalCountsResponse)
async def api_get_total_counts():
    """Get total success and failed request counts."""
    success, failed = await stats_db.get_total_counts()
    return TotalCountsResponse(success=success, failed=failed)


# ---------------------------------------------------------------------------
# OpenAI-compatible API
# ---------------------------------------------------------------------------

def _extract_bearer_token(authorization: Optional[str]) -> str:
    if not authorization:
        raise HTTPException(status_code=401, detail="Missing Authorization header")

    value = authorization.strip()
    if value.lower().startswith("bearer "):
        value = value[7:].strip()
    if not value:
        raise HTTPException(status_code=401, detail="Missing Authorization token")
    return value


async def _resolve_openai_access_token(authorization: Optional[str]) -> str:
    token = _extract_bearer_token(authorization)
    valid_keys = [key.strip() for key in OPENAI_API_KEY.split(",") if key.strip()]
    if not valid_keys:
        return token

    if token not in valid_keys:
        raise HTTPException(status_code=401, detail="Invalid API Key")

    accounts = await list_accounts()
    if not accounts:
        raise HTTPException(status_code=503, detail="No FigureLabs account available")
    return accounts[0]["access_token"]


def _resolve_openai_model_id(model: str) -> int:
    if model in _OPENAI_MODEL_IDS:
        return _OPENAI_MODEL_IDS[model]
    if model.isdigit():
        return int(model)
    raise HTTPException(
        status_code=404,
        detail=f"Model '{model}' not found. Available models: {list(_OPENAI_MODEL_IDS)}",
    )


def _message_content_to_text(content: str | list[dict[str, Any]]) -> str:
    if isinstance(content, str):
        return content

    parts: list[str] = []
    for item in content:
        if item.get("type") == "text" and item.get("text"):
            parts.append(str(item["text"]))
        elif item.get("text"):
            parts.append(str(item["text"]))
    return "\n".join(parts)


def _build_openai_prompt(messages: list[OpenAIMessage]) -> str:
    if not messages:
        raise HTTPException(status_code=400, detail="messages must not be empty")

    parts = [
        (message.role, _message_content_to_text(message.content))
        for message in messages
        if _message_content_to_text(message.content).strip()
    ]
    if not parts:
        raise HTTPException(status_code=400, detail="messages must include text content")
    if len(parts) == 1:
        return parts[0][1]
    return "\n".join(f"{role}: {content}" for role, content in parts)


def _build_figurelabs_result_text(
    session_id: str,
    message_id: str,
    status: Optional[dict[str, Any]],
) -> str:
    lines = [
        "FigureLabs generation completed." if status else "FigureLabs generation submitted.",
        f"session_id: {session_id}",
        f"message_id: {message_id}",
    ]
    urls = (status or {}).get("fileUrl") or []
    for url in urls:
        lines.append(f"file_url: {url}")
    return "\n".join(lines)


async def _run_openai_completion(
    body: OpenAIChatRequest,
    access_token: str,
) -> tuple[str, str, str]:
    prompt = _build_openai_prompt(body.messages)
    model_id = _resolve_openai_model_id(body.model)
    chat = FigureLabsChat(access_token)

    session_id = await asyncio.to_thread(chat.create_session, prompt[:80] or "OpenAI Chat", 0)
    if not session_id:
        raise HTTPException(status_code=502, detail="Failed to create FigureLabs session")

    message_id = await asyncio.to_thread(
        chat.send_message,
        session_id,
        prompt,
        "NORMAL_CHAT",
        model_id,
        body.ratio,
        body.style,
        True,
        body.scene,
    )
    if not message_id:
        raise HTTPException(status_code=502, detail="Failed to send FigureLabs message")

    status = await asyncio.to_thread(chat.wait_for_completion, message_id, body.wait_timeout, 3)
    return _build_figurelabs_result_text(session_id, message_id, status), session_id, message_id


def _chat_completion_response(
    completion_id: str,
    created: int,
    model: str,
    content: str,
) -> dict[str, Any]:
    return {
        "id": completion_id,
        "object": "chat.completion",
        "created": created,
        "model": model,
        "choices": [
            {
                "index": 0,
                "message": {"role": "assistant", "content": content},
                "logprobs": None,
                "finish_reason": "stop",
            }
        ],
        "usage": {"prompt_tokens": 0, "completion_tokens": 0, "total_tokens": 0},
    }


def _chat_completion_chunk(
    completion_id: str,
    created: int,
    model: str,
    delta: dict[str, Any],
    finish_reason: Optional[str] = None,
) -> str:
    return json.dumps(
        {
            "id": completion_id,
            "object": "chat.completion.chunk",
            "created": created,
            "model": model,
            "choices": [
                {
                    "index": 0,
                    "delta": delta,
                    "logprobs": None,
                    "finish_reason": finish_reason,
                }
            ],
        },
        ensure_ascii=False,
    )


@app.get("/v1/models")
async def openai_list_models():
    now = int(time.time())
    return {
        "object": "list",
        "data": [
            {
                "id": model,
                "object": "model",
                "created": now,
                "owned_by": "figurelabs",
                "permission": [],
            }
            for model in _OPENAI_MODEL_IDS
        ],
    }


@app.get("/v1/models/{model_id}")
async def openai_get_model(model_id: str):
    _resolve_openai_model_id(model_id)
    return {"id": model_id, "object": "model", "owned_by": "figurelabs"}


@app.post("/v1/chat/completions")
async def openai_chat_completions(
    body: OpenAIChatRequest,
    authorization: Optional[str] = Header(None),
):
    access_token = await _resolve_openai_access_token(authorization)
    created = int(time.time())
    completion_id = f"chatcmpl-{created}"

    if body.stream:
        async def stream_events():
            role_chunk = _chat_completion_chunk(
                completion_id,
                created,
                body.model,
                {"role": "assistant"},
            )
            yield f"data: {role_chunk}\n\n"

            content, _, message_id = await _run_openai_completion(body, access_token)
            chunk_id = f"chatcmpl-{message_id}"
            content_chunk = _chat_completion_chunk(
                chunk_id,
                created,
                body.model,
                {"content": content},
            )
            yield f"data: {content_chunk}\n\n"

            stop_chunk = _chat_completion_chunk(chunk_id, created, body.model, {}, "stop")
            yield f"data: {stop_chunk}\n\n"
            yield "data: [DONE]\n\n"

        return StreamingResponse(stream_events(), media_type="text/event-stream")

    content, _, message_id = await _run_openai_completion(body, access_token)
    return _chat_completion_response(f"chatcmpl-{message_id}", created, body.model, content)


# ---------------------------------------------------------------------------
# Static frontend
# ---------------------------------------------------------------------------

_STATIC = Path(__file__).parent.parent.parent / "frontend" / "dist"
if _STATIC.exists():
    app.mount("/", StaticFiles(directory=str(_STATIC), html=True), name="static")
